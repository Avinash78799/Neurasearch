"""
NeuraSearch – Document Ingestor
Orchestrates the full ingestion pipeline:
    parse → chunk → embed (ChromaDB) → index (BM25)

Supports PDF and plain-text files, as well as raw byte uploads from the
API layer. Workspace isolation is enforced.
"""

from __future__ import annotations

import logging
import os
import tempfile
from pathlib import Path
from typing import Any

from langchain_core.documents import Document
from pypdf import PdfReader

from config import settings
from rag.chunker import chunk_text
from rag.semantic_chunker import semantic_chunk
from rag.vectorstore import add_documents
from rag.bm25_index import rebuild_index
from workspace_service import WorkspaceContext
from core.exceptions import IngestionError

logger = logging.getLogger(__name__)

# File extensions recognised as plain text.
_TEXT_EXTENSIONS: set[str] = {".txt", ".md", ".csv", ".json", ".log", ".rst"}


# PDF ingestion


def ingest_pdf(file_path: str, filename: str) -> list[Document]:
    """Extract text from a PDF and return chunked ``Document`` objects."""
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"PDF file not found: {file_path}")

    reader = PdfReader(str(path))
    all_chunks: list[Document] = []

    for page_num, page in enumerate(reader.pages, start=1):
        page_text = page.extract_text() or ""
        if not page_text.strip():
            logger.debug("Skipping blank page %d in %s", page_num, filename)
            continue

        metadata: dict[str, Any] = {
            "source": filename,
            "page_number": page_num,
        }
        if settings.use_semantic_chunker:
            chunks = semantic_chunk(page_text, metadata)
        else:
            chunks = chunk_text(page_text, metadata)
        all_chunks.extend(chunks)

    if not all_chunks:
        raise IngestionError(
            f"No extractable text found in PDF: {filename}"
        )

    logger.info(
        "Parsed PDF %s: %d pages → %d chunks",
        filename,
        len(reader.pages),
        len(all_chunks),
    )
    return all_chunks


# Plain-text ingestion


def ingest_txt(file_path: str, filename: str) -> list[Document]:
    """Read a plain-text file and return chunked ``Document`` objects."""
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"Text file not found: {file_path}")

    text = path.read_text(encoding="utf-8", errors="replace")
    if not text.strip():
        raise IngestionError(f"Text file is empty: {filename}")

    metadata: dict[str, Any] = {
        "source": filename,
        "page_number": 1,
    }
    if settings.use_semantic_chunker:
        chunks = semantic_chunk(text, metadata)
    else:
        chunks = chunk_text(text, metadata)

    logger.info(
        "Parsed text file %s: %d chars → %d chunks",
        filename,
        len(text),
        len(chunks),
    )
    return chunks


def ingest_docx(file_path: str, filename: str) -> list[Document]:
    """Extract text from a DOCX Word document and return chunked Document objects."""
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"DOCX file not found: {file_path}")

    import docx
    doc = docx.Document(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    
    # Also extract table text
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                paragraphs.append(row_text)

    full_text = "\n\n".join(paragraphs)
    if not full_text.strip():
        raise IngestionError(f"No extractable text found in DOCX: {filename}")

    metadata: dict[str, Any] = {
        "source": filename,
        "page_number": 1,
    }
    if settings.use_semantic_chunker:
        chunks = semantic_chunk(full_text, metadata)
    else:
        chunks = chunk_text(full_text, metadata)

    logger.info("Parsed DOCX %s: %d paragraphs → %d chunks", filename, len(paragraphs), len(chunks))
    return chunks


def ingest_pptx(file_path: str, filename: str) -> list[Document]:
    """Extract text from a PPTX PowerPoint presentation and return chunked Document objects."""
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"PPTX file not found: {file_path}")

    from pptx import Presentation
    prs = Presentation(str(path))
    all_chunks: list[Document] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
        
        slide_content = "\n".join(slide_texts)
        if not slide_content.strip():
            continue

        metadata: dict[str, Any] = {
            "source": filename,
            "page_number": slide_idx,
        }
        if settings.use_semantic_chunker:
            chunks = semantic_chunk(slide_content, metadata)
        else:
            chunks = chunk_text(slide_content, metadata)
        all_chunks.extend(chunks)

    if not all_chunks:
        raise IngestionError(f"No extractable text found in PPTX: {filename}")

    logger.info("Parsed PPTX %s: %d slides → %d chunks", filename, len(prs.slides), len(all_chunks))
    return all_chunks


# Orchestrator


def ingest_file(file_path: str, filename: str, context: WorkspaceContext | str | None = None) -> dict[str, Any]:
    """Full ingestion pipeline for a file on disk under a specific workspace."""
    ext = Path(filename).suffix.lower()

    try:
        if ext == ".pdf":
            reader = PdfReader(file_path)
            num_pages = len(reader.pages)
            chunks = ingest_pdf(file_path, filename)
        elif ext in (".docx", ".doc"):
            num_pages = 1
            chunks = ingest_docx(file_path, filename)
        elif ext in (".pptx", ".ppt"):
            from pptx import Presentation
            prs = Presentation(file_path)
            num_pages = len(prs.slides)
            chunks = ingest_pptx(file_path, filename)
        elif ext in _TEXT_EXTENSIONS:
            num_pages = 1
            chunks = ingest_txt(file_path, filename)
        else:
            # Fall back to plain-text ingestion for unknown extensions.
            logger.warning(
                "Unrecognised extension '%s' for %s – treating as plain text",
                ext,
                filename,
            )
            num_pages = 1
            chunks = ingest_txt(file_path, filename)


        # Store in vector DB with workspace.
        add_documents(chunks, context)

        # Keep BM25 in sync for workspace.
        rebuild_index(context)

        stats: dict[str, Any] = {
            "filename": filename,
            "num_chunks": len(chunks),
            "num_pages": num_pages,
            "status": "success",
        }
        logger.info("Ingestion complete under workspace context %s: %s", context, stats)
        return stats

    except Exception as exc:
        logger.exception("Ingestion failed for %s", filename)
        return {
            "filename": filename,
            "num_chunks": 0,
            "num_pages": 0,
            "status": f"error: {exc}",
        }


# Bytes-based ingestion (API uploads)


def ingest_bytes(content: bytes, filename: str, context: WorkspaceContext | str | None = None) -> dict[str, Any]:
    """Ingest a document provided as raw bytes under a specific workspace."""
    ext = Path(filename).suffix
    tmp_fd, tmp_path = tempfile.mkstemp(suffix=ext)

    try:
        with os.fdopen(tmp_fd, "wb") as fh:
            fh.write(content)

        return ingest_file(tmp_path, filename, context)
    finally:
        # Clean up temp file regardless of outcome.
        try:
            os.unlink(tmp_path)
        except OSError:
            logger.debug("Could not remove temp file %s", tmp_path)
